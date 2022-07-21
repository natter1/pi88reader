"""
@author: Nathanael Jöhrmann
"""
import os
import statistics
from typing import Union, Iterable, Optional, List

import matplotlib.pyplot as plt
import pptx_tools.style_sheets as style_sheets
from matplotlib.figure import Figure
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx_tools.creator import PPTXCreator, PPTXPosition
from pptx_tools.table_style import PPTXTableStyle
from pptx_tools.templates import analyze_pptx

from pi88reader.ni_analyser import calc_unloading_data, get_power_law_fit_curve
from pi88reader.pi88_importer import PI88Measurement, load_tdm_files
from pi88reader.pi88_plotter import PI88Plotter
from pi88reader.plotter_styles import GraphStyler, get_power_law_fit_curve_style
from pi88reader.pptx_styles import table_style_summary, table_style_measurements_meta
from pi88reader.utils_pi88measurement import get_measurement_result_table_data, get_measurement_meta_table_data
from pi88reader.utils_pi88measurements import get_measurements_meta_table_data, \
    get_measurements_result_table_data


class PI88ToPPTX:
    def __init__(self, measurements_path=None, template=None):
        self.path = measurements_path
        self.measurements = []
        if self.path:
            self.load_tdm_files(measurements_path)
        self.plotter = PI88Plotter(self.measurements)
        self.pptx_creator = PPTXCreator(template=template)
        self.prs = self.pptx_creator.prs
        self.position = self.pptx_creator.default_position

        self.poisson_ratio = 0.3
        self.beta = 1.0

        self.measurements_unloading_data: dict = {}

    def load_tdm_files(self, path: str, sort_key=os.path.getctime):  # sorted by creation time (using windows)
        self.measurements.extend(load_tdm_files(path, sort_key))

    def add_measurements(self, measurements: Union[PI88Measurement, Iterable[PI88Measurement]]) -> None:
        """
        Adds a single PI88Measurement or a list o PI88Measurement's to the plotter.
        """
        if measurements:
            try:
                self.measurements.extend(measurements)
            except TypeError:
                self.measurements.append(measurements)

    def clear_measurements(self):
        self.measurements = []
        self.measurements_unloading_data = {}

    def add_matplotlib_figure(self, fig: Figure, slide: Slide, position: PPTXPosition = None, **kwargs):
        """
        :param fig:
        :param slide_index:
        :param position:
        :param kwargs: e.g. width and height
        :return: prs.shapes.picture.Picture
        """
        return self.pptx_creator.add_matplotlib_figure(fig, slide, position, **kwargs)

    def create_summary_slide(self, title: str = None, layout=None):
        if title is None:
            title = f"Summary - {self.path}"
        result = self.pptx_creator.add_slide(title, layout)

        plotter = PI88Plotter(self.measurements)
        fig = plotter.get_load_displacement_plot()
        fig.axes[0].legend(loc="best")
        self.add_matplotlib_figure(fig, result, PPTXPosition(0.02, 0.15))
        self.create_measurements_result_data_table(result)
        return result

    def create_modulus_hardness_summary_slide(self, layout=None):
        title = "Summary - reduced modulus and hardness"
        result = self.pptx_creator.add_slide(title, layout)

        plotter = PI88Plotter(self.measurements)
        fig = plotter.get_reduced_modulus_plot()
        self.add_matplotlib_figure(fig, result, PPTXPosition(0.02, 0.15))
        avg_reduced_modulus = statistics.mean(fig.axes[0].lines[0].get_ydata())

        fig = plotter.get_hardness_plot()
        self.add_matplotlib_figure(fig, result, PPTXPosition(0.52, 0.15))
        avg_hardness = statistics.mean(fig.axes[0].lines[0].get_ydata())

        text = f"avg. Er = {avg_reduced_modulus} GPa   -   avg. H = {avg_hardness}"
        self.pptx_creator.add_text_box(result, text, PPTXPosition(0.05, 0.85))

        return result

    def create_title_slide(self, title=None, layout=None, default_content=False):
        if title is None:
            title = f"NI results {self.path}"
        result = self.pptx_creator.add_title_slide(title, layout)
        self.create_measurements_meta_table(result)

        plotter = PI88Plotter(self.measurements)
        fig = plotter.get_load_displacement_plot()
        self.add_matplotlib_figure(fig, result, PPTXPosition(0.57, 0.24))
        return result

    def create_measurement_slide(self, measurement: PI88Measurement, layout = None, graph_styler = None):
        title = measurement._name  # filename[:-4].split("/")[-1].split("\\")[-1]
        result = self.pptx_creator.add_slide(title, layout)

        self.create_measurement_result_table(result, measurement)
        self.create_measurement_meta_data_table(result, measurement)

        plotter = PI88Plotter(measurement)
        if graph_styler is not None:
            plotter.graph_styler = graph_styler
        fig = plotter.get_load_displacement_plot()

        if (measurement, self.poisson_ratio, self.beta) in self.measurements_unloading_data:
            fit_data = self.measurements_unloading_data[(measurement, self.poisson_ratio, self.beta)]
            fit_disp, fit_load = get_power_law_fit_curve(**fit_data)
            fig.axes[0].plot(fit_disp, fit_load, **get_power_law_fit_curve_style().dict, label="power-law-fit")

        self.add_matplotlib_figure(fig, result, PPTXPosition(0.02, 0.15))
        return result

    def create_measurement_slides(self, measurements: Optional[List[PI88Measurement]] = None, layout = None) -> list:
        result = []
        if measurements is None:
            measurements = self.measurements

        graph_styler = GraphStyler(len(self.measurements))

        for measurement in measurements:
            result.append(self.create_measurement_slide(measurement, layout, graph_styler))

        return result

    def create_measurement_meta_data_table(self, slide, measurement, table_style: PPTXTableStyle = None) -> Shape:
        table_data = get_measurement_meta_table_data(measurement)
        result = self.pptx_creator.add_table(slide, table_data, PPTXPosition(0.521, 0.16))
        if table_style is None:
            table_style = style_sheets.table_no_header()
            table_style.set_width_as_fraction(0.4)
        table_style.write_shape(result)
        return result

    def _get_measurement_result_table_data(self, measurement: PI88Measurement, poisson_ratio: float, beta: float) -> list:
        if (measurement, poisson_ratio, beta) in self.measurements_unloading_data:
            data = self.measurements_unloading_data[(measurement, poisson_ratio, beta)]
        else:
            data = calc_unloading_data(measurement, beta=self.beta, poisson_ratio=self.poisson_ratio)
            self.measurements_unloading_data[(measurement, poisson_ratio, beta)] = data
        return get_measurement_result_table_data(measurement, data)

    def create_measurement_result_table(self, slide, measurement, table_style: PPTXTableStyle = None) -> Shape:
        table_data = self._get_measurement_result_table_data(measurement, self.poisson_ratio, self.beta)

        result = self.pptx_creator.add_table(slide, table_data, PPTXPosition(0.521, 0.56))
        if table_style is None:
            table_style = style_sheets.table_no_header()
            table_style.set_width_as_fraction(0.4)
        table_style.write_shape(result)
        return result

    def create_measurements_meta_table(self, slide, table_style: PPTXTableStyle = None):
        table_data = get_measurements_meta_table_data(self.measurements)

        result = self.pptx_creator.add_table(slide, table_data)
        if table_style is None:
            table_style = table_style_measurements_meta()
        table_style.write_shape(result)
        return result

    def create_measurements_result_data_table(self, slide, table_style: PPTXTableStyle = None):
        for measurement in self.measurements:
            if measurement not in self.measurements_unloading_data.keys():
                self._get_measurement_result_table_data(measurement, self.poisson_ratio, self.beta)  # todo: better method (name change?)
        table_data = get_measurements_result_table_data(self.measurements_unloading_data.values())
        result = self.pptx_creator.add_table(slide, table_data)
        if table_style is None:
            table_style = table_style_summary()
        table_style.write_shape(result)
        return result

    def get_average_measurements_unloading_data(self):
        sum_Er = sum_E = sum_hardness = 0
        for data in self.measurements_unloading_data.values():
            sum_Er += data['Er']
            sum_E += data['E']
            sum_hardness += data['hardness']
        n = len(self.measurements_unloading_data.values())
        if n > 0:
            return sum_Er/n, sum_E/n, sum_hardness/n
        else:
            return 0, 0, 0
    def save(self, filename="delme.prs"):
        self.prs.save(filename)
