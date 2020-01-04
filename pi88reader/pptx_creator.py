"""
This module provides an easier Interface to create *.pptx presentations using the module python-pptx.
@author: Nathanael Jöhrmann
"""
import io

from pptx import Presentation
from pptx.util import Inches

import pi88reader.pptx_template as pptx_template


class PPTXPosition:
    """
    Used to generate positions of elements in slide coordiinates
    """
    def __init__(self, presentation, left_rel=0.0, top_rel=0.0, left=0, top=0):
        """
        :param presentation: pptx.prs (needed for slide width and height)
        :param left_rel: distance from slide left (relative to slide width)
        :param top_rel: distance from slide top (relative to slide height)
        :param left: "left" to position figure [inches] starting from rel_left
        :param top: "top" to position figure [inches] starting from rel_top
        """
        self.prs = presentation
        self.left_rel = left_rel
        self.top_rel = top_rel
        self.left = left
        self.top = top

    def set(self, left_rel=0.0, top_rel=0.0, left=0, top=0):
        self.left_rel = left_rel
        self.top_rel = top_rel
        self.left = left
        self.top = top

    def dict_for_position(self, left_rel=0.0, top_rel=0.0, left=0, top=0):
        """
        Returns kwargs dict for given position. Does not change attributes of self
        :param left_rel: float [slide_width]
        :param top_rel: float [slide_height]
        :param left: float [inch]
        :param top: float [inch]
        :return: dictionary
        """
        left = self.fraction_width_to_inch(left_rel) + left
        top = self.fraction_height_to_inch(top_rel) + top
        return {"left": left, "top": top}

    def dict(self):
        """
        This method returns a kwargs dict containing "left" and "top".
        :return: dictionary
        """
        return self.dict_for_position(self.left_rel, self.top_rel, self.left, self.top)

    def fraction_width_to_inch(self, fraction):
        """
        Returns a width in inches calculated as a fraction of total slide-width.
        :param fraction: float
        :return: Calculated Width in inch
        """
        result = Inches(self.prs.slide_width.inches * fraction)
        return result

    def fraction_height_to_inch(self, fraction):
        """
        Returns a height in inches calculated as a fraction of total slide-height.
        :param fraction: float
        :return: Calculated Width in inch
        """
        return Inches(self.prs.slide_height.inches * fraction)


# todo: template_file to Enum in pptx_template with all available templates
class PPTXCreator:
    """
    This Class provides an easy interface to create a PowerPoint presentation.
        - position elements as fraction of slide height/width
        - add matplotlib figures
        - use pptx templates (in combination with pptx_template.py)
    """
    def __init__(self, template=None):
        """
        :param template:
        :param title:
        """
        self.slides = []
        self.template = None
        self.prs = None
        self.create_presentation(template)
        self.title_layout = self.prs.slide_masters[0].slide_layouts[0]
        self.default_layout = self.prs.slide_masters[1].slide_layouts[0]
        self.position = PPTXPosition(self.prs)
        # self.template_file = template_file
        # self.set_first_slide(title=title)

        # picture = self.add_matplot_figure(fig, self.prs.slides[0], width=Inches(fig_height*zoom))

        # picture.left = Inches(1)
        # picture.top = Inches(3)

    def fraction_width_to_inch(self, fraction):
        """
        Returns a width in inches calculated as a fraction of total slide-width.
        :param fraction: float
        :return: Calculated Width in inch
        """
        result = Inches(self.prs.slide_width.inches * fraction)
        return result

    def fraction_height_to_inch(self, fraction):
        """
        Returns a height in inches calculated as a fraction of total slide-height.
        :param fraction: float
        :return: Calculated Width in inch
        """
        return Inches(self.prs.slide_height.inches * fraction)

    def save(self, filename="delme.pptx"):
        """
        Saves the presentation under the given filename.
        :param filename: string
        :return: None
        """
        self.prs.save(filename)

    def create_presentation(self, template=None):
        """

        :param template:
        :return:
        """
        if template:
            self.template = template  # pptx_template.TemplateETIT169()
            self.prs = self.template.prs
        else:
            self.prs = Presentation()

    def create_title_slide(self, title, layout=None):
        if not layout:
            layout = self.title_layout
        self.add_slide(title, layout)

    def add_slide(self, title, layout=None):
        if not layout:
            layout = self.default_layout
        slide = self.prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        self.remove_unpopulated_shapes(slide)
        return slide

    def write_position_in_kwargs(self, left_rel=0.0, top_rel=0.0, kwargs={}):
        """
        This method modifies(!) the argument kwargs by adding or changing the entries "left" and "top".
        :param left_rel:
        :param top_rel:
        :param kwargs:
        :return: None
        """
        left = self.fraction_width_to_inch(left_rel)
        top = self.fraction_height_to_inch(top_rel)

        if not "left" in kwargs:
            kwargs["left"] = left
        else:
            kwargs["left"] = kwargs["left"] + left
        if not "top" in kwargs:
            kwargs["top"] = top
        else:
            kwargs["top"] = kwargs["top"] + top

    def add_matplotlib_figure(self, fig, slide_index, pptx_position, **kwargs):
        """
        Add a motplotlib figure fig to slide with index slide_index. With top_rel and left_rel
        it is possible to position the figure in Units of slide height/width (float in range [0, 1].
        :param pptx_position: PPTXPosition
        :param fig: a matplolib figure
        :param slide_index: index of slide in presentation on which to insert fig
        :param kwargs:
        :return: pptx.shapes.picture.Picture
        """
        if not pptx_position:
            pptx_position = self.position
        kwargs.update(pptx_position.dict())
        # todo: check slide_index
        slide = self.prs.slides[slide_index]
        with io.BytesIO() as output:
            fig.savefig(output, format="png")
            pic = slide.shapes.add_picture(output, **kwargs)  # 0, 0)#, left, top)
        return pic

    @staticmethod
    def remove_unpopulated_shapes(slide):
        """
        Removes empty placeholders (e.g. due to layout) from slide.
        Further testing needed.
        :param slide: pptx.slide.Slide
        :return:
        """
        for index in reversed(range(len(slide.shapes))):
            shape = slide.shapes[index]
            # if shape.is_placeholder and shape.text_frame.text == "":
            if shape.has_text_frame and shape.text_frame.text == "":
                shape.element.getparent().remove(shape.element)
                print(f"removed index {index}")
