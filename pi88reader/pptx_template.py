"""
This file contains variables with names of important pptx template master_slide shapes
"""
# from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation
from datetime import datetime


def analyze_pptx(template_file):
    """ Take the given file and analyze the structure of master slides.
    Prints shape names/ids and texts for SlideMaster-shapes
    To get an output file contains marked up information
    remove comment on last two lines of function.
    This is helpful when manipulating template-files.
    """
    prs = Presentation(template_file)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    slide_masters = prs.slide_masters
    for index, slide_master in enumerate(prs.slide_masters):
        print('------------------------------------')
        print('------------------------------------')
        print(f"slide master indexed: {index}")
        print(slide_master)
        print("text boxes:")
        for shape in slide_master.shapes:
            try:
                dummystring = f"shape name: {shape.name} - shape text: {shape.text}"
                shape.text = shape.name
                print(dummystring)
            except:
                pass
            #shape.text = 'hahahaha'
        # for shape in slide_master.slideshapes:
        #     print(shape)
        print('------------------------------------')
        for index, slide_layout in enumerate(slide_master.slide_layouts):
            print(f"\tslide layout: {slide_layout.name}")
            slide = prs.slides.add_slide(slide_master.slide_layouts[index])
            # Not every slide has to have a title
            try:
                title = slide.shapes.title
                title.text = 'Title for Layout {}'.format(index)
            except AttributeError:
                print("No Title for Layout {}".format(index))
            # Go through all the placeholders and identify them by index and type
            for shape in slide.placeholders:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    # Do not overwrite the title which is just a special placeholder
                    try:
                        if 'Title' not in shape.text:
                            shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                    except AttributeError:
                        print("{} has no text attribute".format(phf.type))
                    print(f"\t\tid: {phf.idx} - name: {shape.name}")
    # output_file = '..\\resources\pptx\\template_names.pptx'
    # prs.save(output_file)


# todo: add something like this directly to python-pptx
def copy_font(_from, _to):
    _to.bold = _from.bold
    # todo: color is ColorFormat object
    # _to.set_color = _from.color
    # todo: fill is FillFormat object
    # _to.fill = _from.fill
    _to.italic = _from.italic
    _to.language_id = _from.language_id
    _to.name = _from.name
    _to.size = _from.size
    _to.underline = _from.underline


def remove_unpopulated_shapes(slide):
    """
    Removes empty placeholders (e.g. due to layout) from slide.
    Further testing needed.
    :param slide: pptx.slide.Slide
    :return:
    """
    for index in reversed(range(len(slide.shapes))):
        shape = slide.shapes[index]
        # if shape.is_placeholder and shape.text_frame.text == "":
        if shape.has_text_frame and shape.text_frame.text == "":
            shape.element.getparent().remove(shape.element)
            print(f"removed index {index}")


def change_paragraph_text_to(paragraph, text):
    """
    Change text of paragraph to text, but keep format.
    :param paragraph:
    :param text:
    :return:
    """
    font = paragraph.runs[0].font
    paragraph.text = text
    copy_font(font, paragraph.font)


def analyze_paragraphs(paragraphs):
    for index, para in enumerate(paragraphs):
        print(f"index: {index} - text: {para.text}")
        for run_index, run in enumerate(para.runs):
            print(f"\trun: {run.text}")


# ----------------------------------------------------------------------------
# --------- Customized template classes are needed for each template ---------
# ----------------------------------------------------------------------------
class TemplateETIT169:
    """
    Class handling ETIT 16:9 template.
    """
    TEMPLATE_FILE = '..\\resources\pptx\\ETIT_16-9.pptx'

    def __init__(self):
        self.prs = Presentation(self.TEMPLATE_FILE)

        self.slide_master_big = self.prs.slide_masters[0]
        self.slide_master_small = self.prs.slide_masters[1]

        self.big_layouts = {}
        self.small_layouts = {}
        # following names are the same for small and big master
        self.author_shape_name = "Rectangle 4"
        self.website_shape_name = "Rectangle 5"

        date_time = datetime.now().strftime("%d %B, %Y")
        self.set_author("Nathanael Jöhrmann", city="Chemnitz", date=date_time)
        self.set_website("https://www.tu-chemnitz.de/etit/wetel/")

    def set_author(self, name, city=None, date=None):
        text= ""
        spacer = " ∙ "
        if city:
            text += city + spacer
        if date:
            text += date + spacer
        text += name
        self.write_text_to_master_shape(text=text, shape_name=self.author_shape_name)

    def set_website(self, text):
        self.write_text_to_master_shape(text = text, shape_name=self.website_shape_name)

    def write_text_to_master_shape(self, text, shape_name):
        for shape in self.master_shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == shape_name:
                change_paragraph_text_to(shape.text_frame.paragraphs[0], text)
    @property
    def master_shapes(self):
        result = []
        result.extend(self.slide_master_big.shapes)
        result.extend(self.slide_master_small.shapes)
        return result
